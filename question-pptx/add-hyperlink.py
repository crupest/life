import pptx
import sys
import os

fn = sys.argv[1]

presentation: pptx.presentation.Presentation = pptx.Presentation(fn)

target_slide = presentation.slides.add_slide(presentation.slide_layouts[6])

slide: pptx.slide.Slide
for slide in presentation.slides:
    for shape in slide.shapes:
        try:
            if isinstance(shape, pptx.shapes.autoshape.Shape) and shape.auto_shape_type == pptx.enum.shapes.MSO_SHAPE.OVAL:
                shape.click_action.target_slide = target_slide
        except:
            pass
try:
    os.mkdir('output')
except:
    pass

presentation.save(os.path.join('output', os.path.basename(fn)))
