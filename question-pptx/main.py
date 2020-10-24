import random
import pptx
from pptx.util import Pt
import re
from typing import Iterable

# Chinese colon and English colon is different but both should be supported with love!
choice_question_regex = re.compile(
    r'(?P<problem_id>\d+)\s*[\.、．]?\s*(?P<problem_1>.*?)[\(（]\s*(?P<answer>[a-eA-E])\s*[\)）](?P<problem_2>.*?)\s*[aA]\s*[\.、．]?\s*(?P<choice_a>\S+?)\s*[bB]\s*[\.、．]?\s*(?P<choice_b>\S+?)\s*[cC]\s*[\.、．]?\s*(?P<choice_c>\S+?)\s+([dD]\s*[\.、．]?\s*(?P<choice_d>\S+?)\s+)?([eE]\s*[\.、．]?\s*(?P<choice_e>\S+?)\s+)?', re.MULTILINE)

choice_text = open('./real-choose.txt', 'r', encoding='utf-8').read()

questions = []

for match in choice_question_regex.finditer(choice_text):
    q = {
        'id': int(match.group('problem_id')),
        'type': 'c',
        'problem': match.group('problem_1') + "（ ）" + match.group('problem_2'),
        'answer': match.group('answer').lower(),
        'choices': {
            'a': match.group('choice_a'),
            'b': match.group('choice_b'),
            'c': match.group('choice_c')
        }
    }

    if match.group('choice_d') != None:
        q['choices']['d'] = match.group('choice_d')

    if match.group('choice_e') != None:
        q['choices']['e'] = match.group('choice_e')

    questions.append(q)

print(len(questions))

judge_question_regex = re.compile(
    r'(?P<problem_id>\d+)\s*[\.、．]?\s*(?P<problem>.+)\s*[\(（]\s*(?P<answer>[✓√×xX])\s*[\)）]', re.MULTILINE)

judge_text = open('./real-judge.txt', 'r', encoding='utf-8').read()


for match in judge_question_regex.finditer(judge_text):
    questions.append({
        'id': int(match.group('problem_id')),
        'type': 'j',
        'problem': match.group('problem'),
        'answer': '×' if match.group('answer') in '×xX' else '✓'
    })

print(len(questions))

id_list = [q['id'] for q in questions]
id_list.sort()
for i in range(606):
    if not i in id_list:
        print(i)

random.shuffle(questions)

presentation: pptx.presentation.Presentation = pptx.Presentation()

buttons: Iterable[pptx.shapes.autoshape.Shape] = []

for question in questions:
    slide: pptx.slide.Slide = presentation.slides.add_slide(
        presentation.slide_layouts[6])

    title = slide.shapes.add_textbox(Pt(20), Pt(40), Pt(680), Pt(50))
    problem = slide.shapes.add_textbox(Pt(20), Pt(120), Pt(680), Pt(350))
    answer = slide.shapes.add_textbox(Pt(560), Pt(450), Pt(40), Pt(40))
    id = slide.shapes.add_textbox(Pt(20), Pt(480), Pt(80), Pt(40))

    if question['type'] == 'c':
        title.text = '选择题'
        problem.text = question['problem'] + '\n' + '\n'.join(
            [choice.upper() + '. ' + question['choices'][choice]for choice in question['choices'].keys()])
        answer.text = question['answer'].upper()
    else:
        title.text = '判断题'
        problem.text = question['problem']
        answer.text = question['answer']

    id.text = str(question['id'])

    title_text_frame: pptx.text.text.TextFrame = title.text_frame
    para: pptx.text.text._Paragraph
    for para in title_text_frame.paragraphs:
        para.font.size = Pt(42)
        para.font.color.theme_color = pptx.enum.dml.MSO_THEME_COLOR.TEXT_2

    problem_text_frame: pptx.text.text.TextFrame = problem.text_frame
    problem_text_frame.margin_top = Pt(10)
    para: pptx.text.text._Paragraph
    for para in problem_text_frame.paragraphs:
        para.space_after = Pt(6)
        para.font.size = Pt(30)
    problem_text_frame.word_wrap = True

    answer_text_frame: pptx.text.text.TextFrame = answer.text_frame
    para: pptx.text.text._Paragraph
    for para in answer_text_frame.paragraphs:
        para.font.size = Pt(40)
        para.font.color.rgb = pptx.dml.color.RGBColor(255, 0, 0)

    id_text_frame: pptx.text.text.TextFrame = id.text_frame
    para: pptx.text.text._Paragraph
    for para in id_text_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = pptx.dml.color.RGBColor(128, 128, 128)

    buttons.append(slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.OVAL, Pt(
        660), Pt(460), Pt(40), Pt(40)))

slide: pptx.slide.Slide = presentation.slides.add_slide(
    presentation.slide_layouts[6])

for button in buttons:
    button.click_action.target_slide = slide

presentation.save('test.pptx')
